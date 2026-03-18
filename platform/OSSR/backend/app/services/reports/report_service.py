"""
OSSR Research Report Service (AntiGravity Agent — S7)
Generates analytical reports from research simulations and topic data.
Two report types: Research Evolution and Comparative Field.
"""

import json
import threading
from typing import Dict, Any, List, Optional
from dataclasses import dataclass, field
from datetime import datetime

from opensens_common.config import Config
from ...models.research import ResearchDataStore, TopicLevel
from opensens_common.task import TaskManager, TaskStatus
from opensens_common.llm_client import LLMClient
from ..agents.profile_gen import ResearcherProfileStore
from ..simulation.runner import ResearchSimulationRunner
from ...db import get_connection

import logging

logger = logging.getLogger(__name__)


# ── Report Types ─────────────────────────────────────────────────────


REPORT_TYPES = {
    "evolution": {
        "name": "Research Evolution Report",
        "description": (
            "Analyzes how a research topic evolves through agent discourse. "
            "Tracks position shifts, emerging consensus, and dissenting views across rounds."
        ),
    },
    "comparative": {
        "name": "Comparative Field Report",
        "description": (
            "Compares two or more research subfields, identifying overlap, gaps, "
            "potential synergies, and cross-pollination opportunities."
        ),
    },
    "infographic": {
        "name": "Visual Infographic",
        "description": (
            "Structured data for visual infographic rendering of key findings, "
            "including stats, flow diagrams, and color-coded topic relationships."
        ),
    },
}


# ── LLM Prompts ──────────────────────────────────────────────────────


EVOLUTION_PLAN_PROMPT = """\
You are a research analyst generating a report on how academic discourse evolved \
during a simulated research discussion.

**Simulation context:**
- Topic: {topic}
- Format: {format_name}
- Participants: {agent_count} researcher agents
- Rounds completed: {rounds}

**Agent roster:**
{agent_roster}

**Discussion transcript (summarized):**
{transcript_summary}

Generate a JSON report outline with 3-5 sections. Focus on:
1. How positions evolved across rounds
2. Key arguments and counter-arguments
3. Emerging consensus or persistent disagreements
4. Novel ideas or synthesis that arose
5. Implications for future research directions

Output JSON:
{{
    "title": "Report title",
    "summary": "One-line executive summary",
    "sections": [
        {{
            "title": "Section title",
            "description": "What this section covers",
            "key_points": ["point 1", "point 2"]
        }}
    ]
}}"""


COMPARATIVE_PLAN_PROMPT = """\
You are a research analyst generating a comparative report across research subfields.

**Fields being compared:**
{fields_list}

**Topic landscape summary:**
- Total papers: {paper_count}
- Topic clusters: {cluster_count}
- Identified gaps: {gap_count}

**Field details:**
{field_details}

**Research gaps between fields:**
{gaps_summary}

Generate a JSON report outline with 3-5 sections. Focus on:
1. Strengths and focus areas of each field
2. Methodological differences and complementarities
3. Citation overlap and cross-pollination patterns
4. Gap opportunities for interdisciplinary work
5. Predicted convergence or divergence trajectories

Output JSON:
{{
    "title": "Report title",
    "summary": "One-line executive summary",
    "sections": [
        {{
            "title": "Section title",
            "description": "What this section covers",
            "key_points": ["point 1", "point 2"]
        }}
    ]
}}"""


SECTION_WRITE_PROMPT = """\
You are writing a section of a research analysis report.

**Report:** {report_title}
**Section:** {section_title}
**Section focus:** {section_description}
**Key points to cover:** {key_points}

**Context data:**
{context_data}

Write this section in clear, academic English. Requirements:
- Use specific evidence from the provided data (cite DOIs where available)
- Include direct quotes from agent discussions where relevant (use > blockquote format)
- Be analytical, not merely descriptive
- 400-800 words per section
- Use markdown formatting (headers, bold, lists where appropriate)

Write the section content now:"""


# ── Data Classes ─────────────────────────────────────────────────────


@dataclass
class ReportSection:
    title: str
    description: str = ""
    key_points: List[str] = field(default_factory=list)
    content: str = ""

    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "description": self.description,
            "key_points": self.key_points,
            "content": self.content,
        }


@dataclass
class ResearchReport:
    report_id: str
    report_type: str  # "evolution" or "comparative"
    title: str = ""
    summary: str = ""
    sections: List[ReportSection] = field(default_factory=list)
    status: str = "pending"  # pending, generating, completed, failed
    simulation_id: Optional[str] = None
    topic_ids: List[str] = field(default_factory=list)
    created_at: str = field(default_factory=lambda: datetime.now().isoformat())
    error: Optional[str] = None

    def to_dict(self) -> Dict[str, Any]:
        return {
            "report_id": self.report_id,
            "report_type": self.report_type,
            "title": self.title,
            "summary": self.summary,
            "sections": [s.to_dict() for s in self.sections],
            "status": self.status,
            "simulation_id": self.simulation_id,
            "topic_ids": self.topic_ids,
            "created_at": self.created_at,
            "error": self.error,
        }

    def to_markdown(self) -> str:
        lines = [f"# {self.title}", "", f"*{self.summary}*", ""]
        for section in self.sections:
            lines.append(f"## {section.title}")
            lines.append("")
            lines.append(section.content)
            lines.append("")
        return "\n".join(lines)

    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'ResearchReport':
        sections = [
            ReportSection(
                title=s["title"],
                description=s.get("description", ""),
                key_points=s.get("key_points", []),
                content=s.get("content", ""),
            )
            for s in data.get("sections", [])
        ]
        return cls(
            report_id=data["report_id"],
            report_type=data["report_type"],
            title=data.get("title", ""),
            summary=data.get("summary", ""),
            sections=sections,
            status=data.get("status", "pending"),
            simulation_id=data.get("simulation_id"),
            topic_ids=data.get("topic_ids", []),
            created_at=data.get("created_at", ""),
            error=data.get("error"),
        )


# ── Report Store ─────────────────────────────────────────────────────


class ResearchReportStore:
    """SQLite-backed store for generated reports. Singleton."""

    _instance = None

    def __new__(cls):
        if cls._instance is None:
            cls._instance = super().__new__(cls)
        return cls._instance

    def add(self, report: ResearchReport):
        conn = get_connection()
        conn.execute(
            "INSERT OR REPLACE INTO reports (report_id, data) VALUES (?, ?)",
            (report.report_id, json.dumps(report.to_dict())),
        )
        conn.commit()

    def get(self, report_id: str) -> Optional[ResearchReport]:
        conn = get_connection()
        row = conn.execute(
            "SELECT data FROM reports WHERE report_id = ?", (report_id,)
        ).fetchone()
        if row is None:
            return None
        return ResearchReport.from_dict(json.loads(row["data"]))

    def list_all(self) -> List[ResearchReport]:
        conn = get_connection()
        rows = conn.execute("SELECT data FROM reports").fetchall()
        return [ResearchReport.from_dict(json.loads(r["data"])) for r in rows]


# ── Report Generator ─────────────────────────────────────────────────


class ResearchReportGenerator:
    """
    Generates research reports from simulation transcripts or topic landscape data.
    """

    def __init__(self):
        self.data_store = ResearchDataStore()
        self.profile_store = ResearcherProfileStore()
        self.sim_runner = ResearchSimulationRunner()
        self.report_store = ResearchReportStore()
        self.task_manager = TaskManager()

    # ── Public API ────────────────────────────────────────────────

    def generate_evolution_report(self, simulation_id: str) -> str:
        """
        Generate a research evolution report from a completed simulation.
        Returns task_id for async tracking.
        """
        import uuid

        report_id = f"rpt_{uuid.uuid4().hex[:12]}"
        report = ResearchReport(
            report_id=report_id,
            report_type="evolution",
            simulation_id=simulation_id,
        )
        self.report_store.add(report)

        task_id = self.task_manager.create_task(
            task_type="research_report",
            metadata={"report_id": report_id, "type": "evolution"},
        )

        thread = threading.Thread(
            target=self._evolution_worker,
            args=(task_id, report),
            daemon=True,
        )
        thread.start()
        return task_id

    def generate_comparative_report(self, topic_ids: List[str]) -> str:
        """
        Generate a comparative field report from topic landscape data.
        Returns task_id for async tracking.
        """
        import uuid

        report_id = f"rpt_{uuid.uuid4().hex[:12]}"
        report = ResearchReport(
            report_id=report_id,
            report_type="comparative",
            topic_ids=topic_ids,
        )
        self.report_store.add(report)

        task_id = self.task_manager.create_task(
            task_type="research_report",
            metadata={"report_id": report_id, "type": "comparative"},
        )

        thread = threading.Thread(
            target=self._comparative_worker,
            args=(task_id, report),
            daemon=True,
        )
        thread.start()
        return task_id

    def get_report(self, report_id: str) -> Optional[ResearchReport]:
        return self.report_store.get(report_id)

    def list_reports(self) -> List[ResearchReport]:
        return self.report_store.list_all()

    # ── Export Methods ─────────────────────────────────────────────

    def export_pptx(self, report_id: str) -> bytes:
        """Export a completed report as a PowerPoint (.pptx) file."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        import io

        report = self.report_store.get(report_id)
        if not report:
            raise ValueError(f"Report not found: {report_id}")
        if report.status != "completed":
            raise ValueError(f"Report not completed (status={report.status})")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = report.title
        if len(slide.placeholders) > 1 and slide.placeholders[1]:
            slide.placeholders[1].text = report.summary

        # One content slide per section
        content_layout = prs.slide_layouts[1]  # Title and Content
        for section in report.sections:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = section.title

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            # Add key points as bullets
            if section.key_points:
                for i, kp in enumerate(section.key_points):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = kp
                    p.font.size = Pt(16)
                    p.level = 0

            # Add content excerpt
            if section.content:
                # Strip markdown formatting for plain text
                plain = section.content.replace("**", "").replace("###", "").replace("> ", "")
                excerpt = plain[:800] + ("..." if len(plain) > 800 else "")
                p = tf.add_paragraph()
                p.text = ""
                p = tf.add_paragraph()
                p.text = excerpt
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.LEFT

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        buf.close()
        return data

    def export_audio(self, report_id: str) -> bytes:
        """Export a TTS audio summary of the report using OpenAI TTS API."""
        import openai
        import os

        report = self.report_store.get(report_id)
        if not report:
            raise ValueError(f"Report not found: {report_id}")
        if report.status != "completed":
            raise ValueError(f"Report not completed (status={report.status})")

        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            raise ValueError("OPENAI_API_KEY not configured — TTS export unavailable")

        # Build condensed text for TTS (max ~4000 chars for reasonable audio length)
        parts = [report.title, report.summary]
        for section in report.sections:
            parts.append(f"Section: {section.title}")
            if section.key_points:
                parts.append(". ".join(section.key_points))
            if section.content:
                # Take first 500 chars of content per section
                plain = section.content.replace("**", "").replace("###", "").replace("> ", "")
                parts.append(plain[:500])

        text = ". ".join(parts)
        if len(text) > 4000:
            text = text[:3997] + "..."

        try:
            client = openai.OpenAI(api_key=api_key)
            response = client.audio.speech.create(
                model="tts-1",
                voice="alloy",
                input=text,
            )
            return response.content
        except Exception as e:
            logger.error(f"OpenAI TTS API call failed: {e}")
            raise ValueError(f"TTS generation failed: {e}")

    def generate_infographic(self, report_id: str) -> Dict[str, Any]:
        """Generate structured infographic data from a report using Gemini LLM."""
        report = self.report_store.get(report_id)
        if not report:
            raise ValueError(f"Report not found: {report_id}")
        if report.status != "completed":
            raise ValueError(f"Report not completed (status={report.status})")

        # Build report digest for the LLM
        digest_parts = [f"Title: {report.title}", f"Summary: {report.summary}"]
        for section in report.sections:
            digest_parts.append(f"\nSection: {section.title}")
            if section.key_points:
                digest_parts.append(f"Key points: {', '.join(section.key_points)}")
            if section.content:
                digest_parts.append(section.content[:600])
        digest = "\n".join(digest_parts)

        prompt = (
            "You are a data visualization expert. Given the following research report, "
            "generate a structured JSON infographic specification that a frontend can render.\n\n"
            f"Report:\n{digest[:4000]}\n\n"
            "Output JSON with this structure:\n"
            '{\n'
            '  "title": "Infographic title",\n'
            '  "subtitle": "One-line summary",\n'
            '  "color_scheme": ["#hex1", "#hex2", "#hex3", "#hex4"],\n'
            '  "stats": [{"label": "...", "value": "...", "icon": "chart|people|paper|globe"}],\n'
            '  "flow_nodes": [{"id": "n1", "label": "...", "description": "...", "connects_to": ["n2"]}],\n'
            '  "key_findings": ["finding 1", "finding 2", "finding 3"],\n'
            '  "sections": [{"title": "...", "bullets": ["..."]}]\n'
            '}\n\n'
            "Generate 3-5 stats, 4-8 flow nodes showing concept relationships, "
            "and 3-5 key findings. Use an academic color scheme."
        )

        # Try Gemini first, fall back to any available provider
        llm = None
        for provider in ["gemini", "openai", "anthropic"]:
            try:
                llm = LLMClient(provider=provider)
                logger.info(f"Infographic LLM initialized: provider={provider}")
                break
            except (ValueError, KeyError) as e:
                logger.debug(f"LLMClient init skipped for {provider}: {e}")
                continue

        if llm:
            try:
                result = llm.chat_json(
                    messages=[
                        {"role": "system", "content": "You are a data visualization expert."},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.4,
                )
                result["report_id"] = report_id
                result["report_title"] = report.title
                return result
            except Exception as e:
                logger.warning(f"LLM infographic generation failed: {e}")

        # Fallback: build basic infographic data from report structure
        return {
            "report_id": report_id,
            "report_title": report.title,
            "title": report.title,
            "subtitle": report.summary,
            "color_scheme": ["#1e3a5f", "#2d6a4f", "#e76f51", "#264653"],
            "stats": [
                {"label": "Sections", "value": str(len(report.sections)), "icon": "chart"},
                {"label": "Key Points", "value": str(sum(len(s.key_points) for s in report.sections)), "icon": "paper"},
            ],
            "flow_nodes": [
                {"id": f"s{i}", "label": s.title, "description": s.description, "connects_to": [f"s{i+1}"] if i < len(report.sections) - 1 else []}
                for i, s in enumerate(report.sections)
            ],
            "key_findings": [kp for s in report.sections for kp in s.key_points[:2]][:5],
            "sections": [{"title": s.title, "bullets": s.key_points} for s in report.sections],
        }

    # ── Evolution Report Worker ───────────────────────────────────

    def _evolution_worker(self, task_id: str, report: ResearchReport):
        self.task_manager.update_task(
            task_id, status=TaskStatus.PROCESSING, progress=0,
            message="Preparing evolution report...",
        )
        report.status = "generating"

        try:
            # Gather simulation data
            sim = self.sim_runner.get_simulation(report.simulation_id)
            if not sim:
                raise ValueError(f"Simulation not found: {report.simulation_id}")

            transcript = self.sim_runner.get_transcript(report.simulation_id)

            # Build agent roster
            agent_lines = []
            for aid in sim.agent_ids:
                profile = self.profile_store.get(aid)
                if profile:
                    agent_lines.append(
                        f"- {profile.name} ({profile.role}, {profile.primary_field}): {profile.bio}"
                    )
                else:
                    agent_lines.append(f"- {aid}")

            # Summarize transcript (take first ~3000 chars worth)
            transcript_summary = self._summarize_transcript(transcript)

            # Step 1: Plan outline via LLM
            self.task_manager.update_task(
                task_id, progress=10, message="Planning report outline..."
            )

            outline = self._plan_evolution_outline(
                topic=sim.topic,
                format_name=sim.discussion_format,
                agent_count=len(sim.agent_ids),
                rounds=sim.current_round,
                agent_roster="\n".join(agent_lines),
                transcript_summary=transcript_summary,
            )

            report.title = outline["title"]
            report.summary = outline["summary"]
            report.sections = [
                ReportSection(
                    title=s["title"],
                    description=s.get("description", ""),
                    key_points=s.get("key_points", []),
                )
                for s in outline["sections"]
            ]

            # Step 2: Generate each section
            total_sections = len(report.sections)
            for i, section in enumerate(report.sections):
                pct = 20 + int(70 * i / max(total_sections, 1))
                self.task_manager.update_task(
                    task_id, progress=pct,
                    message=f"Writing section {i+1}/{total_sections}: {section.title}",
                )

                context = self._build_evolution_context(sim, transcript, i)
                section.content = self._write_section(report, section, context)

            # Done
            report.status = "completed"
            self.report_store.add(report)
            self.task_manager.complete_task(task_id, result={
                "report_id": report.report_id,
                "title": report.title,
                "sections": total_sections,
            })

        except Exception as e:
            logger.exception(f"Evolution report failed: {e}")
            report.status = "failed"
            report.error = str(e)
            self.report_store.add(report)
            self.task_manager.fail_task(task_id, str(e))

    # ── Comparative Report Worker ─────────────────────────────────

    def _comparative_worker(self, task_id: str, report: ResearchReport):
        self.task_manager.update_task(
            task_id, status=TaskStatus.PROCESSING, progress=0,
            message="Preparing comparative report...",
        )
        report.status = "generating"

        try:
            # Gather topic data
            topics = []
            for tid in report.topic_ids:
                t = self.data_store.get_topic(tid)
                if t:
                    topics.append(t)

            if len(topics) < 2:
                raise ValueError("At least 2 topics required for comparative report")

            # Build field details
            field_details_lines = []
            for t in topics:
                paper_ids = self.data_store.get_topic_papers(t.topic_id)
                papers = []
                for pid in paper_ids:
                    p = self.data_store.get_paper_by_id(pid)
                    if p:
                        papers.append(p)

                kw_set = set()
                for p in papers:
                    kw_set.update(p.keywords[:5])

                field_details_lines.append(
                    f"**{t.name}** (Level {t.level.value}): "
                    f"{len(papers)} papers, "
                    f"keywords: {', '.join(list(kw_set)[:10])}"
                )

            # Get gaps
            from .research_mapper import ResearchMapper
            mapper = ResearchMapper()
            all_gaps = mapper.find_gaps(min_score=0.2)
            topic_names = {t.name for t in topics}
            relevant_gaps = [
                g for g in all_gaps
                if g.get("topic_a") in topic_names or g.get("partner_topic") in topic_names
            ]

            gaps_summary = ""
            for g in relevant_gaps[:10]:
                gaps_summary += (
                    f"- {g['topic_a']} <-> {g['partner_topic']}: "
                    f"score {g.get('gap_score', 0):.2f}"
                )
                if g.get("opportunity"):
                    gaps_summary += f" — {g['opportunity']}"
                gaps_summary += "\n"

            if not gaps_summary:
                gaps_summary = "(No specific gaps identified between these fields)"

            stats = self.data_store.stats()

            # Step 1: Plan outline
            self.task_manager.update_task(
                task_id, progress=10, message="Planning comparative outline..."
            )

            outline = self._plan_comparative_outline(
                fields_list=", ".join(t.name for t in topics),
                paper_count=stats.get("total_papers", 0),
                cluster_count=stats.get("total_topics", 0),
                gap_count=len(relevant_gaps),
                field_details="\n".join(field_details_lines),
                gaps_summary=gaps_summary,
            )

            report.title = outline["title"]
            report.summary = outline["summary"]
            report.sections = [
                ReportSection(
                    title=s["title"],
                    description=s.get("description", ""),
                    key_points=s.get("key_points", []),
                )
                for s in outline["sections"]
            ]

            # Step 2: Generate each section
            total_sections = len(report.sections)
            for i, section in enumerate(report.sections):
                pct = 20 + int(70 * i / max(total_sections, 1))
                self.task_manager.update_task(
                    task_id, progress=pct,
                    message=f"Writing section {i+1}/{total_sections}: {section.title}",
                )

                context = self._build_comparative_context(topics, relevant_gaps, i)
                section.content = self._write_section(report, section, context)

            report.status = "completed"
            self.report_store.add(report)
            self.task_manager.complete_task(task_id, result={
                "report_id": report.report_id,
                "title": report.title,
                "sections": total_sections,
            })

        except Exception as e:
            logger.exception(f"Comparative report failed: {e}")
            report.status = "failed"
            report.error = str(e)
            self.report_store.add(report)
            self.task_manager.fail_task(task_id, str(e))

    # ── LLM Helpers ───────────────────────────────────────────────

    def _get_llm(self) -> Optional[LLMClient]:
        try:
            return LLMClient()
        except ValueError:
            return None

    def _plan_evolution_outline(self, **kwargs) -> Dict[str, Any]:
        llm = self._get_llm()
        prompt = EVOLUTION_PLAN_PROMPT.format(**kwargs)

        if llm:
            try:
                return llm.chat_json(
                    messages=[
                        {"role": "system", "content": "You are a research report planner."},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.4,
                )
            except Exception as e:
                logger.warning(f"LLM outline planning failed: {e}")

        # Fallback
        return {
            "title": f"Research Evolution: {kwargs.get('topic', 'Unknown')}",
            "summary": "Analysis of research discourse evolution across simulation rounds.",
            "sections": [
                {
                    "title": "Opening Positions",
                    "description": "Initial stances and perspectives of participating researchers.",
                    "key_points": ["Starting positions", "Methodological preferences"],
                },
                {
                    "title": "Discourse Dynamics",
                    "description": "How arguments developed and shifted through the discussion.",
                    "key_points": ["Key debates", "Position shifts", "Citation patterns"],
                },
                {
                    "title": "Outcomes and Directions",
                    "description": "Emerging consensus, unresolved tensions, and future directions.",
                    "key_points": ["Consensus areas", "Open questions", "Recommended directions"],
                },
            ],
        }

    def _plan_comparative_outline(self, **kwargs) -> Dict[str, Any]:
        llm = self._get_llm()
        prompt = COMPARATIVE_PLAN_PROMPT.format(**kwargs)

        if llm:
            try:
                return llm.chat_json(
                    messages=[
                        {"role": "system", "content": "You are a research report planner."},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.4,
                )
            except Exception as e:
                logger.warning(f"LLM outline planning failed: {e}")

        fields = kwargs.get("fields_list", "Field A, Field B")
        return {
            "title": f"Comparative Analysis: {fields}",
            "summary": f"Cross-field comparison of {fields}.",
            "sections": [
                {
                    "title": "Field Profiles",
                    "description": "Overview of each field's focus, methods, and key contributions.",
                    "key_points": ["Research focus", "Dominant methodologies", "Publication volume"],
                },
                {
                    "title": "Overlap and Divergence",
                    "description": "Where fields intersect and where they diverge.",
                    "key_points": ["Shared keywords", "Citation overlap", "Methodological differences"],
                },
                {
                    "title": "Gap Opportunities",
                    "description": "Identified gaps and cross-pollination potential.",
                    "key_points": ["Underexplored intersections", "Synergy potential", "Suggested directions"],
                },
            ],
        }

    def _write_section(
        self,
        report: ResearchReport,
        section: ReportSection,
        context_data: str,
    ) -> str:
        llm = self._get_llm()
        if not llm:
            # Fallback: structured summary from key points
            lines = [f"### {section.title}", ""]
            if section.description:
                lines.append(section.description)
                lines.append("")
            for kp in section.key_points:
                lines.append(f"- **{kp}**")
            lines.append("")
            lines.append("*Detailed analysis requires LLM configuration.*")
            return "\n".join(lines)

        prompt = SECTION_WRITE_PROMPT.format(
            report_title=report.title,
            section_title=section.title,
            section_description=section.description,
            key_points=", ".join(section.key_points) if section.key_points else "(none specified)",
            context_data=context_data,
        )

        try:
            result = llm.chat(
                messages=[
                    {"role": "system", "content": "You are a research analyst writing a report section."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.5,
            )
            return result.strip()
        except Exception as e:
            logger.warning(f"Section generation failed: {e}")
            return f"*Section generation failed: {e}*"

    # ── Context Builders ──────────────────────────────────────────

    def _summarize_transcript(self, transcript: List[Dict], max_chars: int = 3000) -> str:
        if not transcript:
            return "(No transcript available)"

        lines = []
        total = 0
        for turn in transcript:
            line = (
                f"[Round {turn.get('round_num', '?')}] "
                f"{turn.get('agent_name', 'Unknown')} ({turn.get('role', '')}): "
                f"{turn.get('content', '')[:200]}"
            )
            if total + len(line) > max_chars:
                lines.append("... (transcript truncated)")
                break
            lines.append(line)
            total += len(line)
        return "\n".join(lines)

    def _build_evolution_context(
        self,
        sim,
        transcript: List[Dict],
        section_index: int,
    ) -> str:
        """Build context data for an evolution report section."""
        total_rounds = sim.current_round or 1

        if section_index == 0:
            # First section: focus on early rounds
            relevant = [t for t in transcript if t.get("round_num", 0) <= max(1, total_rounds // 3)]
        elif section_index == len(sim.agent_ids) - 1 or section_index >= 2:
            # Last section: focus on later rounds
            cutoff = max(1, total_rounds * 2 // 3)
            relevant = [t for t in transcript if t.get("round_num", 0) >= cutoff]
        else:
            # Middle sections: full transcript
            relevant = transcript

        lines = []
        for turn in relevant[:30]:
            content = turn.get("content", "")[:300]
            cited = turn.get("cited_dois", [])
            cite_str = f" [cites: {', '.join(cited)}]" if cited else ""
            lines.append(
                f"**{turn.get('agent_name', '?')}** (Round {turn.get('round_num', '?')}, "
                f"{turn.get('role', '')}): {content}{cite_str}"
            )
        return "\n\n".join(lines) if lines else "(No relevant transcript data)"

    def _build_comparative_context(
        self,
        topics: list,
        gaps: List[Dict],
        section_index: int,
    ) -> str:
        """Build context data for a comparative report section."""
        lines = []

        for t in topics:
            paper_ids = self.data_store.get_topic_papers(t.topic_id)
            papers = []
            for pid in paper_ids[:15]:
                p = self.data_store.get_paper_by_id(pid)
                if p:
                    papers.append(p)

            lines.append(f"### {t.name}")
            lines.append(f"Papers: {len(paper_ids)}")

            for p in papers[:5]:
                kw = ", ".join(p.keywords[:4])
                lines.append(f"- {p.title} ({p.source}) — keywords: {kw}")
            lines.append("")

        if gaps and section_index >= 1:
            lines.append("### Identified Gaps")
            for g in gaps[:8]:
                lines.append(
                    f"- {g.get('topic_a', '?')} <-> {g.get('partner_topic', '?')}: "
                    f"score={g.get('gap_score', 0):.2f}"
                )
                if g.get("opportunity"):
                    lines.append(f"  Opportunity: {g['opportunity']}")

        return "\n".join(lines) if lines else "(No topic data available)"
